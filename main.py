import os
import copy
import six
from pptx import Presentation

input_dir_files = os.listdir('.')
powerpoints = []
for x in input_dir_files: 
    if x.endswith(".pptx"):
        powerpoints.append(x)
        
output_prs = Presentation()

for prs in powerpoints:
    for slide in prs:
        print("hek")




output_prs.save("new-presentation.pptx")

#https://github.com/scanny/python-pptx/issues/132


 ## potential fixes
def duplicate_slide(pres, index):
    template = pres.slides[index]
    try:
        blank_slide_layout = pres.slide_layouts[12]
    except:
        blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(
                value.reltype,
                value._target,
                value.rId
            )

    return copied_slide


def _get_blank_slide_layout(pres):
         layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
         min_items = min(layout_items_count)
         blank_layout_id = layout_items_count.index(min_items)
         return pres.slide_layouts[blank_layout_id]

    def copy_slide(pres,pres1,index):
         source = pres.slides[index]

         blank_slide_layout = _get_blank_slide_layout(pres)
         dest = pres1.slides.add_slide(blank_slide_layout)

         for shp in source.shapes:
              el = shp.element
              newel = copy.deepcopy(el)
              dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

              for key, value in six.iteritems(source.rels):
                         # Make sure we don't copy a notesSlide relation as that won't exist
                       if not "notesSlide" in value.reltype:
                               dest.rels.add_relationship(value.reltype, value._target, value.rId)

              return dest


def duplicate_slide2(pres, index):
    """Duplicate the slide with the given index in pres.

    Adds slide to the end of the presentation"""
    source = pres.slides[index]

    blank_slide_layout = _get_blank_slide_layout(pres)
    dest = pres.slides.add_slide(blank_slide_layout)

    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for key, value in six.iteritems(source.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if not "notesSlide" in value.reltype:
            dest.rels.add_relationship(value.reltype, value._target, value.rId)

    return dest